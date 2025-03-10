from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List
from datetime import datetime
from io import BytesIO
from pptx import Presentation
import logging

router = APIRouter()
logger = logging.getLogger(__name__)

class PPTContent(BaseModel):
    title: str
    agenda: str
    key_findings: List[str]
    recommendations: List[str]
    conclusion: str
    qa_points: List[str]

@router.post("/generate")
async def generate_ppt(ppt_content: PPTContent):
    try:
        current_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # Create PPT
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = ppt_content.title
        title_slide.placeholders[1].text = f"Generated on {current_time}"
        
        # Agenda slide
        agenda_slide = prs.slides.add_slide(prs.slide_layouts[1])
        agenda_slide.shapes.title.text = "Agenda"
        agenda_slide.placeholders[1].text = ppt_content.agenda
        
        # Key Findings
        findings_slide = prs.slides.add_slide(prs.slide_layouts[1])
        findings_slide.shapes.title.text = "Key Findings"
        findings_text = findings_slide.placeholders[1].text_frame
        for finding in ppt_content.key_findings:
            p = findings_text.add_paragraph()
            p.text = f"• {finding}"
        
        # Recommendations
        recom_slide = prs.slides.add_slide(prs.slide_layouts[1])
        recom_slide.shapes.title.text = "Recommendations"
        recom_text = recom_slide.placeholders[1].text_frame
        for rec in ppt_content.recommendations:
            p = recom_text.add_paragraph()
            p.text = f"• {rec}"
        
        # Conclusion
        conclusion_slide = prs.slides.add_slide(prs.slide_layouts[1])
        conclusion_slide.shapes.title.text = "Conclusion"
        conclusion_slide.placeholders[1].text = ppt_content.conclusion
        
        # Q&A
        qa_slide = prs.slides.add_slide(prs.slide_layouts[1])
        qa_slide.shapes.title.text = "Q&A Discussion Points"
        qa_text = qa_slide.placeholders[1].text_frame
        for qa in ppt_content.qa_points:
            p = qa_text.add_paragraph()
            p.text = f"• {qa}"
        
        # Save to BytesIO
        ppt_file = BytesIO()
        prs.save(ppt_file)
        ppt_file.seek(0)

        return StreamingResponse(
            ppt_file,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": 'attachment; filename="AWS_Analysis.pptx"'}
        )

    except Exception as e:
        logger.error(f"PPT Generation Error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))
